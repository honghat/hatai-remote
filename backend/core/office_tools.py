"""
Office Tools — Professional Excel, Word, PowerPoint manipulation for HatAI Agent.

Design principles:
- Simple JSON args that a small LLM can generate reliably
- Each tool does ONE clear thing
- Rich defaults so the model doesn't need to specify every detail
- Returns structured previews so the model can verify results
"""
import os
import json
import csv
import logging
import datetime
from typing import Dict, Any, List, Optional

logger = logging.getLogger("OfficeTools")

ALLOWED_PATHS = ["/Users/nguyenhat", "/Volumes/HatAI", "/tmp"]


def _check_path(path: str) -> bool:
    abs_path = os.path.realpath(path)
    return any(abs_path.startswith(os.path.realpath(b)) for b in ALLOWED_PATHS)


def _ensure_dir(path: str):
    d = os.path.dirname(path)
    if d:
        os.makedirs(d, exist_ok=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_read(args: Dict[str, Any]) -> Dict[str, Any]:
    """Read Excel (.xlsx/.xls) or CSV file. Returns sheet names, headers, and data preview.

    Args:
        path: file path
        sheet: sheet name (default: first sheet, or "all" for all sheets)
        max_rows: max rows to return (default 50)
    """
    import openpyxl

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    max_rows = min(int(args.get("max_rows", 50)), 200)

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    ext = os.path.splitext(path)[1].lower()

    # CSV handling
    if ext in (".csv", ".tsv"):
        try:
            delimiter = "\t" if ext == ".tsv" else ","
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = []
                for i, row in enumerate(reader):
                    if i > max_rows:
                        break
                    rows.append(row)
            headers = rows[0] if rows else []
            data = rows[1:] if len(rows) > 1 else []
            return {
                "path": path,
                "format": ext,
                "headers": headers,
                "data": data,
                "row_count": len(data),
                "col_count": len(headers),
            }
        except Exception as e:
            return {"error": f"Failed to read CSV: {str(e)}"}

    # Excel handling
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames

        sheets_to_read = []
        if sheet_name == "all":
            sheets_to_read = sheet_names
        elif sheet_name and sheet_name in sheet_names:
            sheets_to_read = [sheet_name]
        else:
            sheets_to_read = [sheet_names[0]]

        result = {
            "path": path,
            "format": ext,
            "sheet_names": sheet_names,
            "sheets": {},
        }

        for sn in sheets_to_read:
            ws = wb[sn]
            rows_data = []
            headers = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(c) if c is not None else "" for c in row]
                else:
                    rows_data.append([_cell_to_str(c) for c in row])
                if i >= max_rows:
                    break

            result["sheets"][sn] = {
                "headers": headers,
                "data": rows_data,
                "row_count": ws.max_row or 0,
                "col_count": ws.max_column or 0,
            }

        wb.close()

        # Shortcut: if only 1 sheet, flatten
        if len(sheets_to_read) == 1:
            sn = sheets_to_read[0]
            result["headers"] = result["sheets"][sn]["headers"]
            result["data"] = result["sheets"][sn]["data"]
            result["row_count"] = result["sheets"][sn]["row_count"]
            result["col_count"] = result["sheets"][sn]["col_count"]

        return result

    except Exception as e:
        return {"error": f"Failed to read Excel: {str(e)}"}


def _cell_to_str(val):
    if val is None:
        return ""
    if isinstance(val, datetime.datetime):
        return val.strftime("%Y-%m-%d %H:%M")
    if isinstance(val, datetime.date):
        return val.strftime("%Y-%m-%d")
    return str(val)


def tool_excel_write(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a new Excel file with data, formatting, and formulas.

    Args:
        path: output file path (.xlsx)
        sheets: list of sheet definitions, each with:
            - name: sheet name
            - headers: list of column headers
            - data: list of rows (each row is a list of values)
            - col_widths: optional list of column widths (numbers)
            - formulas: optional list of {cell, formula} e.g. [{"cell":"C2", "formula":"=A2+B2"}]
        style: optional style preset: "default", "professional", "colorful", "minimal"
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    path = args.get("path", "")
    sheets = args.get("sheets", [])
    style_preset = args.get("style", "professional")

    if not path:
        return {"error": "path is required"}
    if not path.endswith(".xlsx"):
        path += ".xlsx"
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}

    # Support simple single-sheet shortcut
    if not sheets:
        headers = args.get("headers", [])
        data = args.get("data", [])
        if headers or data:
            sheets = [{"name": "Sheet1", "headers": headers, "data": data}]

    if not sheets:
        return {"error": "sheets (or headers+data) required"}

    # Style presets
    STYLES = {
        "professional": {
            "header_font": Font(bold=True, color="FFFFFF", size=11),
            "header_fill": PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid"),
            "header_alignment": Alignment(horizontal="center", vertical="center"),
            "alt_fill": PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid"),
            "border": Border(
                left=Side(style="thin", color="B4C6E7"),
                right=Side(style="thin", color="B4C6E7"),
                top=Side(style="thin", color="B4C6E7"),
                bottom=Side(style="thin", color="B4C6E7"),
            ),
        },
        "colorful": {
            "header_font": Font(bold=True, color="FFFFFF", size=11),
            "header_fill": PatternFill(start_color="FF6B35", end_color="FF6B35", fill_type="solid"),
            "header_alignment": Alignment(horizontal="center", vertical="center"),
            "alt_fill": PatternFill(start_color="FFF0E8", end_color="FFF0E8", fill_type="solid"),
            "border": Border(
                left=Side(style="thin", color="FFB088"),
                right=Side(style="thin", color="FFB088"),
                top=Side(style="thin", color="FFB088"),
                bottom=Side(style="thin", color="FFB088"),
            ),
        },
        "minimal": {
            "header_font": Font(bold=True, size=11),
            "header_fill": PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid"),
            "header_alignment": Alignment(horizontal="left"),
            "alt_fill": None,
            "border": Border(bottom=Side(style="thin", color="CCCCCC")),
        },
        "default": {
            "header_font": Font(bold=True),
            "header_fill": None,
            "header_alignment": Alignment(horizontal="center"),
            "alt_fill": None,
            "border": None,
        },
    }
    st = STYLES.get(style_preset, STYLES["professional"])

    try:
        _ensure_dir(path)
        wb = openpyxl.Workbook()
        # Remove default sheet
        wb.remove(wb.active)

        for sheet_def in sheets:
            ws = wb.create_sheet(title=sheet_def.get("name", "Sheet"))
            headers = sheet_def.get("headers", [])
            data = sheet_def.get("data", [])
            col_widths = sheet_def.get("col_widths", [])
            formulas = sheet_def.get("formulas", [])

            # Write headers
            if headers:
                for col_idx, h in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col_idx, value=h)
                    if st["header_font"]:
                        cell.font = st["header_font"]
                    if st["header_fill"]:
                        cell.fill = st["header_fill"]
                    if st["header_alignment"]:
                        cell.alignment = st["header_alignment"]
                    if st["border"]:
                        cell.border = st["border"]

                # Freeze header row
                ws.freeze_panes = "A2"

            # Write data
            start_row = 2 if headers else 1
            for row_idx, row in enumerate(data):
                for col_idx, val in enumerate(row, 1):
                    cell = ws.cell(row=start_row + row_idx, column=col_idx, value=_parse_value(val))
                    if st["border"]:
                        cell.border = st["border"]
                    # Alternating row colors
                    if st["alt_fill"] and row_idx % 2 == 1:
                        cell.fill = st["alt_fill"]

            # Apply formulas
            for f_def in formulas:
                cell_ref = f_def.get("cell", "")
                formula = f_def.get("formula", "")
                if cell_ref and formula:
                    ws[cell_ref] = formula

            # Auto-width or custom widths
            if col_widths:
                for i, w in enumerate(col_widths):
                    ws.column_dimensions[get_column_letter(i + 1)].width = w
            else:
                _auto_width(ws, headers, data)

            # Auto-filter on header row
            if headers:
                ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}{start_row + len(data) - 1}"

        wb.save(path)
        wb.close()

        total_rows = sum(len(s.get("data", [])) for s in sheets)
        return {
            "message": f"Excel created: {path}",
            "path": path,
            "sheets": len(sheets),
            "total_rows": total_rows,
            "style": style_preset,
        }
    except Exception as e:
        return {"error": f"Failed to create Excel: {str(e)}"}


def _parse_value(val):
    """Convert string values to appropriate Python types for Excel."""
    if val is None or val == "":
        return None
    if isinstance(val, (int, float)):
        return val
    s = str(val).strip()
    # Try number
    try:
        if "." in s:
            return float(s.replace(",", ""))
        return int(s.replace(",", ""))
    except ValueError:
        pass
    return s


def _auto_width(ws, headers, data):
    """Auto-fit column widths based on content."""
    from openpyxl.utils import get_column_letter
    max_cols = max(len(headers), max((len(r) for r in data), default=0)) if data else len(headers)
    for col_idx in range(1, max_cols + 1):
        max_len = 0
        letter = get_column_letter(col_idx)
        # Check header
        if col_idx <= len(headers):
            max_len = len(str(headers[col_idx - 1]))
        # Check data (sample first 50 rows)
        for row in data[:50]:
            if col_idx <= len(row):
                max_len = max(max_len, len(str(row[col_idx - 1])))
        ws.column_dimensions[letter].width = min(max(max_len + 3, 8), 50)


def tool_excel_edit(args: Dict[str, Any]) -> Dict[str, Any]:
    """Edit cells in an existing Excel file.

    Args:
        path: file path
        sheet: sheet name (default: active sheet)
        edits: list of cell edits, each with:
            - cell: cell reference like "A1", "B5"
            - value: new value (number, text, or formula starting with =)
        add_rows: optional list of rows to append at the end
        delete_rows: optional list of row numbers to delete (1-indexed)
        add_sheet: optional new sheet name to create
    """
    import openpyxl

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    edits = args.get("edits", [])
    add_rows = args.get("add_rows", [])
    delete_rows = args.get("delete_rows", [])
    add_sheet = args.get("add_sheet", None)

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path)

        # Add new sheet if requested
        if add_sheet:
            wb.create_sheet(title=add_sheet)
            if not sheet_name:
                sheet_name = add_sheet

        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        changes = 0

        # Apply cell edits
        for edit in edits:
            cell_ref = edit.get("cell", "")
            value = edit.get("value", "")
            if cell_ref:
                ws[cell_ref] = _parse_value(value) if not str(value).startswith("=") else value
                changes += 1

        # Append rows
        for row in add_rows:
            ws.append([_parse_value(v) for v in row])
            changes += 1

        # Delete rows (from bottom to top to preserve indices)
        for row_num in sorted(delete_rows, reverse=True):
            ws.delete_rows(row_num)
            changes += 1

        wb.save(path)
        wb.close()

        return {
            "message": f"Excel edited: {changes} changes applied",
            "path": path,
            "sheet": ws.title,
            "changes": changes,
        }
    except Exception as e:
        return {"error": f"Failed to edit Excel: {str(e)}"}


def tool_excel_chart(args: Dict[str, Any]) -> Dict[str, Any]:
    """Add a chart to an Excel file.

    Args:
        path: file path
        sheet: sheet name (default: active)
        chart_type: "bar", "line", "pie", "column", "area", "scatter" (default: column)
        title: chart title
        data_range: e.g. "B1:B10" (Y values)
        categories_range: e.g. "A1:A10" (X labels)
        position: cell to place chart, e.g. "E2" (default: "E2")
        width: chart width in cm (default 18)
        height: chart height in cm (default 12)
    """
    import openpyxl
    from openpyxl.chart import BarChart, LineChart, PieChart, AreaChart, ScatterChart, Reference

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    chart_type = args.get("chart_type", "column")
    title = args.get("title", "Chart")
    data_range = args.get("data_range", "")
    categories_range = args.get("categories_range", "")
    position = args.get("position", "E2")
    width = float(args.get("width", 18))
    height = float(args.get("height", 12))

    if not path or not data_range:
        return {"error": "path and data_range are required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    CHART_CLASSES = {
        "bar": BarChart,
        "column": BarChart,
        "line": LineChart,
        "pie": PieChart,
        "area": AreaChart,
        "scatter": ScatterChart,
    }

    if chart_type not in CHART_CLASSES:
        return {"error": f"Unknown chart_type: {chart_type}. Use: {list(CHART_CLASSES.keys())}"}

    try:
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        chart = CHART_CLASSES[chart_type]()
        chart.title = title
        chart.width = width
        chart.height = height

        if chart_type == "bar":
            chart.type = "bar"
        elif chart_type == "column":
            chart.type = "col"

        # Parse ranges
        data_ref = _parse_range_to_reference(ws, data_range)
        chart.add_data(data_ref, titles_from_data=True)

        if categories_range:
            cat_ref = _parse_range_to_reference(ws, categories_range)
            chart.set_categories(cat_ref)

        # Style
        chart.style = 10

        ws.add_chart(chart, position)
        wb.save(path)
        wb.close()

        return {
            "message": f"{chart_type} chart added to {path}",
            "path": path,
            "chart_type": chart_type,
            "title": title,
            "position": position,
        }
    except Exception as e:
        return {"error": f"Failed to add chart: {str(e)}"}


def _parse_range_to_reference(ws, range_str: str):
    """Convert 'A1:A10' or 'B1:D10' to openpyxl Reference."""
    from openpyxl.chart import Reference
    from openpyxl.utils import column_index_from_string
    import re

    m = re.match(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', range_str.upper())
    if not m:
        raise ValueError(f"Invalid range: {range_str}")

    min_col = column_index_from_string(m.group(1))
    min_row = int(m.group(2))
    max_col = column_index_from_string(m.group(3))
    max_row = int(m.group(4))

    return Reference(ws, min_col=min_col, min_row=min_row, max_col=max_col, max_row=max_row)


def tool_excel_format(args: Dict[str, Any]) -> Dict[str, Any]:
    """Format cells in an existing Excel file.

    Args:
        path: file path
        sheet: sheet name (default: active)
        formats: list of format operations:
            - range: cell range like "A1:D1" or single cell "A1"
            - bold: true/false
            - italic: true/false
            - font_size: number
            - font_color: hex color without # (e.g. "FF0000" for red)
            - bg_color: hex background color (e.g. "FFFF00" for yellow)
            - align: "left", "center", "right"
            - number_format: e.g. "#,##0", "#,##0.00", "0%", "yyyy-mm-dd"
            - merge: true (merge the range)
            - border: "thin", "medium", "thick"
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    formats = args.get("formats", [])

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        changes = 0
        for fmt in formats:
            cell_range = fmt.get("range", "")
            if not cell_range:
                continue

            # Check if merge requested
            if fmt.get("merge"):
                ws.merge_cells(cell_range)
                changes += 1

            # Get cells in range
            cells = _get_cells_in_range(ws, cell_range)

            for cell in cells:
                # Font
                font_kwargs = {}
                if "bold" in fmt:
                    font_kwargs["bold"] = fmt["bold"]
                if "italic" in fmt:
                    font_kwargs["italic"] = fmt["italic"]
                if "font_size" in fmt:
                    font_kwargs["size"] = fmt["font_size"]
                if "font_color" in fmt:
                    font_kwargs["color"] = fmt["font_color"]
                if font_kwargs:
                    # Merge with existing font
                    existing = cell.font
                    cell.font = Font(
                        bold=font_kwargs.get("bold", existing.bold),
                        italic=font_kwargs.get("italic", existing.italic),
                        size=font_kwargs.get("size", existing.size),
                        color=font_kwargs.get("color", existing.color),
                        name=existing.name,
                    )

                # Background
                if "bg_color" in fmt:
                    cell.fill = PatternFill(
                        start_color=fmt["bg_color"],
                        end_color=fmt["bg_color"],
                        fill_type="solid",
                    )

                # Alignment
                if "align" in fmt:
                    cell.alignment = Alignment(horizontal=fmt["align"])

                # Number format
                if "number_format" in fmt:
                    cell.number_format = fmt["number_format"]

                # Border
                if "border" in fmt:
                    side = Side(style=fmt["border"])
                    cell.border = Border(left=side, right=side, top=side, bottom=side)

                changes += 1

        wb.save(path)
        wb.close()

        return {
            "message": f"Formatted {changes} cells in {path}",
            "path": path,
            "changes": changes,
        }
    except Exception as e:
        return {"error": f"Failed to format Excel: {str(e)}"}


def _get_cells_in_range(ws, range_str: str) -> list:
    """Get all cells in a range like 'A1:D5' or single cell 'A1'."""
    if ":" in range_str:
        cells = []
        for row in ws[range_str]:
            if isinstance(row, tuple):
                cells.extend(row)
            else:
                cells.append(row)
        return cells
    else:
        return [ws[range_str]]


def tool_excel_analyze(args: Dict[str, Any]) -> Dict[str, Any]:
    """Quick analysis of Excel data: summary stats, top/bottom values, duplicates.

    Args:
        path: file path
        sheet: sheet name (default: active)
        column: column letter or header name to analyze (e.g. "B" or "Revenue")
        operation: "stats", "unique", "top", "bottom", "duplicates", "count_by"
        n: number of results for top/bottom (default 10)
    """
    import openpyxl

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    column = args.get("column", "")
    operation = args.get("operation", "stats")
    n = int(args.get("n", 10))

    if not path or not column:
        return {"error": "path and column are required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        # Find column index
        col_idx = None
        headers = []
        for i, cell in enumerate(next(ws.iter_rows(min_row=1, max_row=1, values_only=False)), 1):
            h = str(cell.value) if cell.value else ""
            headers.append(h)
            if column.upper() == cell.column_letter or column.lower() == h.lower():
                col_idx = cell.column

        if col_idx is None:
            wb.close()
            return {"error": f"Column '{column}' not found. Headers: {headers}"}

        # Extract values (skip header)
        values = []
        for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx, values_only=True):
            val = row[0]
            if val is not None:
                values.append(val)

        wb.close()

        if not values:
            return {"column": column, "error": "No data in column"}

        # Separate numeric and text values
        nums = [v for v in values if isinstance(v, (int, float))]
        texts = [str(v) for v in values]

        result = {"column": column, "total_values": len(values), "operation": operation}

        if operation == "stats":
            if nums:
                result["numeric_count"] = len(nums)
                result["sum"] = round(sum(nums), 2)
                result["average"] = round(sum(nums) / len(nums), 2)
                result["min"] = min(nums)
                result["max"] = max(nums)
                result["median"] = round(sorted(nums)[len(nums) // 2], 2)
            else:
                result["message"] = "No numeric values, text column"
                result["sample"] = texts[:10]

        elif operation == "unique":
            unique = list(set(texts))
            result["unique_count"] = len(unique)
            result["values"] = unique[:n]

        elif operation == "top":
            if nums:
                sorted_vals = sorted(enumerate(nums), key=lambda x: x[1], reverse=True)[:n]
                result["top"] = [{"row": i + 2, "value": v} for i, v in sorted_vals]
            else:
                result["error"] = "Column is not numeric"

        elif operation == "bottom":
            if nums:
                sorted_vals = sorted(enumerate(nums), key=lambda x: x[1])[:n]
                result["bottom"] = [{"row": i + 2, "value": v} for i, v in sorted_vals]
            else:
                result["error"] = "Column is not numeric"

        elif operation == "duplicates":
            from collections import Counter
            counts = Counter(texts)
            dups = {k: v for k, v in counts.items() if v > 1}
            result["duplicates"] = dict(sorted(dups.items(), key=lambda x: x[1], reverse=True)[:n])
            result["duplicate_count"] = len(dups)

        elif operation == "count_by":
            from collections import Counter
            counts = Counter(texts)
            result["counts"] = dict(sorted(counts.items(), key=lambda x: x[1], reverse=True)[:n])

        return result

    except Exception as e:
        return {"error": f"Failed to analyze: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════════
#  WORD (DOCX) TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_word_read(args: Dict[str, Any]) -> Dict[str, Any]:
    """Read a Word document. Returns paragraphs and tables.

    Args:
        path: file path (.docx)
        max_chars: max characters to return (default 8000)
    """
    from docx import Document

    path = args.get("path", "")
    max_chars = int(args.get("max_chars", 8000))

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        doc = Document(path)

        # Extract paragraphs with style info
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append({
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal",
                })

        # Extract tables
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            tables.append(rows)

        # Build text preview
        text = "\n".join(p["text"] for p in paragraphs)
        if len(text) > max_chars:
            text = text[:max_chars] + "...[truncated]"

        return {
            "path": path,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
            "paragraphs": paragraphs[:100],
            "tables": tables[:10],
            "text_preview": text,
        }
    except Exception as e:
        return {"error": f"Failed to read Word: {str(e)}"}


def tool_word_write(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a professional Word document.

    Args:
        path: output file path (.docx)
        title: document title
        content: list of content blocks, each with:
            - type: "heading", "paragraph", "table", "bullet", "numbered", "pagebreak"
            - text: text content (for heading, paragraph, bullet, numbered)
            - level: heading level 1-4 (for heading, default 1)
            - items: list of strings (for bullet/numbered lists)
            - headers: list of column headers (for table)
            - rows: list of rows (for table)
            - bold: true/false
            - italic: true/false
        style: "professional", "simple" (default: professional)
    """
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    # Robustness: handle case where args is somehow a string
    if isinstance(args, str):
        # If it looks like a path, treat as path
        if args.endswith(".docx") or "/" in args:
            args = {"path": args}
        else:
            return {"error": "Invalid arguments to word_write. Must be a dict."}

    path = args.get("path", "")
    title = args.get("title", "")
    content = args.get("content", [])
    style = args.get("style", "professional")

    if not path:
        return {"error": "path is required"}
    if not path.endswith(".docx"):
        path += ".docx"
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}

    try:
        _ensure_dir(path)
        doc = Document()

        # Set default font
        doc_style = doc.styles["Normal"]
        doc_style.font.name = "Calibri"
        doc_style.font.size = Pt(11)

        # Title
        if title:
            t = doc.add_heading(title, level=0)
            if style == "professional":
                for run in t.runs:
                    run.font.color.rgb = RGBColor(0x2F, 0x54, 0x96)

            # Add date line
            date_para = doc.add_paragraph()
            date_run = date_para.add_run(f"Ngày tạo: {datetime.datetime.now().strftime('%d/%m/%Y')}")
            date_run.font.size = Pt(10)
            date_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            doc.add_paragraph()  # Spacer

        # Content blocks
        for block in content:
            # Handle case where block is a string instead of a dict
            if isinstance(block, str):
                block = {"type": "paragraph", "text": block}
            
            btype = block.get("type", "paragraph")

            if btype == "heading":
                level = min(int(block.get("level", 1)), 4)
                h = doc.add_heading(block.get("text", ""), level=level)
                if style == "professional" and level <= 2:
                    for run in h.runs:
                        run.font.color.rgb = RGBColor(0x2F, 0x54, 0x96)

            elif btype == "paragraph":
                p = doc.add_paragraph()
                text = block.get("text", "")
                run = p.add_run(text)
                if block.get("bold"):
                    run.font.bold = True
                if block.get("italic"):
                    run.font.italic = True

            elif btype == "bullet":
                items = block.get("items", [])
                if isinstance(items, str):
                    items = [items]
                for item in items:
                    doc.add_paragraph(str(item), style="List Bullet")

            elif btype == "numbered":
                items = block.get("items", [])
                if isinstance(items, str):
                    items = [items]
                for item in items:
                    doc.add_paragraph(str(item), style="List Number")

            elif btype == "table":
                headers = block.get("headers", [])
                rows = block.get("rows", [])

                if headers:
                    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                    table.style = "Light Grid Accent 1" if style == "professional" else "Table Grid"
                    table.alignment = WD_TABLE_ALIGNMENT.CENTER

                    # Headers
                    for i, h in enumerate(headers):
                        cell = table.rows[0].cells[i]
                        cell.text = str(h)
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.font.bold = True

                    # Data rows
                    for row_idx, row in enumerate(rows):
                        for col_idx, val in enumerate(row):
                            if col_idx < len(headers):
                                table.rows[row_idx + 1].cells[col_idx].text = str(val)

                    doc.add_paragraph()  # Spacer after table

            elif btype == "pagebreak":
                doc.add_page_break()

        doc.save(path)

        return {
            "message": f"Word document created: {path}",
            "path": path,
            "blocks": len(content),
        }
    except Exception as e:
        return {"error": f"Failed to create Word: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT (PPTX) TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_pptx_read(args: Dict[str, Any]) -> Dict[str, Any]:
    """Read a PowerPoint file. Returns slide titles and content.

    Args:
        path: file path (.pptx)
    """
    from pptx import Presentation

    path = args.get("path", "")

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        prs = Presentation(path)
        slides_info = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            slides_info.append({
                "slide": i,
                "layout": slide.slide_layout.name if slide.slide_layout else "",
                "content": texts,
            })

        return {
            "path": path,
            "slide_count": len(slides_info),
            "slides": slides_info,
        }
    except Exception as e:
        return {"error": f"Failed to read PowerPoint: {str(e)}"}


def tool_pptx_write(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a professional PowerPoint presentation.

    Args:
        path: output file path (.pptx)
        title: presentation title (for title slide)
        subtitle: subtitle for title slide
        slides: list of slide definitions:
            - layout: "title", "content", "two_content", "blank" (default: content)
            - title: slide title
            - body: text content or list of bullet points
            - notes: speaker notes (optional)
            - table: optional {headers, rows} for table slide
        theme_color: hex color for accents (default "2F5496" blue)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = args.get("path", "")
    title = args.get("title", "Presentation")
    subtitle = args.get("subtitle", "")
    slides = args.get("slides", [])
    theme_color = args.get("theme_color", "2F5496")

    if not path:
        return {"error": "path is required"}
    if not path.endswith(".pptx"):
        path += ".pptx"
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}

    try:
        _ensure_dir(path)
        prs = Presentation()

        # Parse theme color
        tc = RGBColor(
            int(theme_color[0:2], 16),
            int(theme_color[2:4], 16),
            int(theme_color[4:6], 16),
        )

        # ── Title slide ──
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle

        # Style title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = tc
            run.font.size = Pt(36)

        # ── Content slides ──
        for slide_def in slides:
            layout_name = slide_def.get("layout", "content")
            slide_title = slide_def.get("title", "")
            body = slide_def.get("body", "")
            notes = slide_def.get("notes", "")
            table_def = slide_def.get("table", None)

            # Choose layout index
            layout_idx = {
                "title": 0,
                "content": 1,  # Title + Content
                "two_content": 3,
                "blank": 6,
            }.get(layout_name, 1)

            try:
                slide_layout = prs.slide_layouts[layout_idx]
            except IndexError:
                slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if slide.shapes.title and slide_title:
                slide.shapes.title.text = slide_title
                for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = tc

            # Body content
            if body and len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                tf = ph.text_frame
                tf.clear()

                if isinstance(body, list):
                    # Bullet points
                    for i, item in enumerate(body):
                        if i == 0:
                            tf.paragraphs[0].text = str(item)
                            tf.paragraphs[0].font.size = Pt(18)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(item)
                            p.font.size = Pt(18)
                            p.space_before = Pt(6)
                else:
                    tf.paragraphs[0].text = str(body)
                    tf.paragraphs[0].font.size = Pt(18)

            # Table
            if table_def:
                headers = table_def.get("headers", [])
                rows = table_def.get("rows", [])
                if headers:
                    n_rows = len(rows) + 1
                    n_cols = len(headers)
                    left = Inches(0.5)
                    top = Inches(2.0)
                    width = Inches(9.0)
                    height = Inches(0.3) * n_rows

                    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

                    # Headers
                    for i, h in enumerate(headers):
                        cell = table.cell(0, i)
                        cell.text = str(h)
                        for para in cell.text_frame.paragraphs:
                            para.font.bold = True
                            para.font.size = Pt(12)
                            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = tc

                    # Data
                    for r_idx, row in enumerate(rows):
                        for c_idx, val in enumerate(row):
                            if c_idx < n_cols:
                                cell = table.cell(r_idx + 1, c_idx)
                                cell.text = str(val)
                                for para in cell.text_frame.paragraphs:
                                    para.font.size = Pt(11)

            # Speaker notes
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)

        prs.save(path)

        return {
            "message": f"PowerPoint created: {path}",
            "path": path,
            "slide_count": 1 + len(slides),
        }
    except Exception as e:
        return {"error": f"Failed to create PowerPoint: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════════
#  CSV/DATA CONVERSION TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_csv_to_excel(args: Dict[str, Any]) -> Dict[str, Any]:
    """Convert CSV/TSV to formatted Excel file.

    Args:
        input_path: source CSV/TSV path
        output_path: destination .xlsx path
        style: "professional", "colorful", "minimal" (default: professional)
    """
    input_path = args.get("input_path", "")
    output_path = args.get("output_path", "")

    if not input_path:
        return {"error": "input_path is required"}

    # Read CSV first
    read_result = tool_excel_read({"path": input_path, "max_rows": 10000})
    if "error" in read_result:
        return read_result

    headers = read_result.get("headers", [])
    data = read_result.get("data", [])

    if not output_path:
        output_path = os.path.splitext(input_path)[0] + ".xlsx"

    return tool_excel_write({
        "path": output_path,
        "headers": headers,
        "data": data,
        "style": args.get("style", "professional"),
    })


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL — ROWS/COLUMNS/SHEETS MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_rows_cols(args: Dict[str, Any]) -> Dict[str, Any]:
    """Insert, delete, hide/show rows or columns in Excel.

    Args:
        path: file path
        sheet: sheet name (default: active)
        operations: list of operations:
          - action: "insert_row" | "delete_row" | "insert_col" | "delete_col"
                    | "hide_row" | "show_row" | "hide_col" | "show_col"
                    | "set_row_height" | "set_col_width"
          - idx: row/column number (1-based for rows, 1-based or letter for cols)
          - count: how many to insert/delete (default 1)
          - size: height (pt) for rows or width for columns
    """
    import openpyxl
    from openpyxl.utils import column_index_from_string

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    operations = args.get("operations", [])

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        done = []

        for op in operations:
            action = op.get("action", "")
            idx = op.get("idx", 1)
            count = int(op.get("count", 1))
            size = op.get("size", None)

            # Normalize column idx: "C" → 3
            if isinstance(idx, str) and idx.isalpha():
                col_idx = column_index_from_string(idx.upper())
            else:
                col_idx = int(idx)

            if action == "insert_row":
                ws.insert_rows(col_idx, amount=count)
            elif action == "delete_row":
                ws.delete_rows(col_idx, amount=count)
            elif action == "insert_col":
                ws.insert_cols(col_idx, amount=count)
            elif action == "delete_col":
                ws.delete_cols(col_idx, amount=count)
            elif action == "hide_row":
                for r in range(col_idx, col_idx + count):
                    ws.row_dimensions[r].hidden = True
            elif action == "show_row":
                for r in range(col_idx, col_idx + count):
                    ws.row_dimensions[r].hidden = False
            elif action == "hide_col":
                from openpyxl.utils import get_column_letter
                for c in range(col_idx, col_idx + count):
                    ws.column_dimensions[get_column_letter(c)].hidden = True
            elif action == "show_col":
                from openpyxl.utils import get_column_letter
                for c in range(col_idx, col_idx + count):
                    ws.column_dimensions[get_column_letter(c)].hidden = False
            elif action == "set_row_height":
                for r in range(col_idx, col_idx + count):
                    ws.row_dimensions[r].height = float(size)
            elif action == "set_col_width":
                from openpyxl.utils import get_column_letter
                for c in range(col_idx, col_idx + count):
                    ws.column_dimensions[get_column_letter(c)].width = float(size)
            else:
                done.append(f"unknown action: {action}")
                continue
            done.append(f"{action} idx={idx} count={count}")

        wb.save(path)
        wb.close()
        return {"message": f"{len(done)} operations done", "operations": done, "path": path}
    except Exception as e:
        return {"error": str(e)}


def tool_excel_sheets(args: Dict[str, Any]) -> Dict[str, Any]:
    """Manage sheets: create, rename, delete, copy, reorder, protect.

    Args:
        path: file path
        operations: list of:
          - action: "create" | "rename" | "delete" | "copy" | "move" | "protect" | "unprotect"
          - sheet: target sheet name
          - new_name: new name (for rename/copy)
          - position: 0-based index (for move/create)
          - password: for protect/unprotect
    """
    import openpyxl

    path = args.get("path", "")
    operations = args.get("operations", [])

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path)
        done = []

        for op in operations:
            action = op.get("action", "")
            sheet = op.get("sheet", "")
            new_name = op.get("new_name", "")
            position = op.get("position", None)
            password = op.get("password", "")

            if action == "create":
                ws = wb.create_sheet(title=sheet or "New Sheet",
                                     index=position if position is not None else len(wb.sheetnames))
                done.append(f"created '{ws.title}'")
            elif action == "rename":
                if sheet in wb.sheetnames:
                    wb[sheet].title = new_name
                    done.append(f"renamed '{sheet}' → '{new_name}'")
                else:
                    done.append(f"sheet not found: {sheet}")
            elif action == "delete":
                if sheet in wb.sheetnames and len(wb.sheetnames) > 1:
                    del wb[sheet]
                    done.append(f"deleted '{sheet}'")
                else:
                    done.append(f"cannot delete '{sheet}' (not found or last sheet)")
            elif action == "copy":
                if sheet in wb.sheetnames:
                    src = wb[sheet]
                    tgt = wb.copy_worksheet(src)
                    if new_name:
                        tgt.title = new_name
                    done.append(f"copied '{sheet}' → '{tgt.title}'")
            elif action == "move":
                if sheet in wb.sheetnames and position is not None:
                    wb.move_sheet(sheet, offset=int(position) - wb.sheetnames.index(sheet))
                    done.append(f"moved '{sheet}' to position {position}")
            elif action == "protect":
                if sheet in wb.sheetnames:
                    wb[sheet].protection.sheet = True
                    if password:
                        wb[sheet].protection.password = password
                    done.append(f"protected '{sheet}'")
            elif action == "unprotect":
                if sheet in wb.sheetnames:
                    wb[sheet].protection.sheet = False
                    done.append(f"unprotected '{sheet}'")

        wb.save(path)
        wb.close()
        return {"message": f"{len(done)} sheet operations done", "operations": done, "sheets": wb.sheetnames if hasattr(wb, 'sheetnames') else []}
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL — FORMULAS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_formula(args: Dict[str, Any]) -> Dict[str, Any]:
    """Write formulas, named ranges, data validation and conditional formatting.

    Args:
        path: file path
        sheet: sheet name
        formulas: list of {cell, formula} — formula starts with =
          e.g. {"cell":"C2","formula":"=SUM(A2:B2)"}
        named_ranges: list of {name, ref} — e.g. {"name":"Revenue","ref":"Sheet1!B2:B100"}
        conditional_formats: list of:
          - range: e.g. "C2:C100"
          - rule: "greater" | "less" | "equal" | "between" | "top10" | "data_bar" | "color_scale"
          - value: threshold value (for greater/less/equal)
          - value2: second value (for between)
          - color: hex color for highlight (e.g. "FF0000")
        data_validation: list of:
          - range: e.g. "D2:D100"
          - type: "list" | "whole" | "decimal" | "date"
          - formula1: for list: "Item1,Item2" or cell range "$G$1:$G$10"; for number: min value
          - formula2: max value (for between)
          - message: optional input message
    """
    import openpyxl
    from openpyxl.formatting.rule import ColorScaleRule, DataBarRule, CellIsRule, Rule
    from openpyxl.styles import PatternFill
    from openpyxl.utils import quote_sheetname

    path = args.get("path", "")
    sheet_name = args.get("sheet", None)
    formulas = args.get("formulas", [])
    named_ranges = args.get("named_ranges", [])
    cond_fmts = args.get("conditional_formats", [])
    validations = args.get("data_validation", [])

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        done = []

        # Write formulas
        for f in formulas:
            cell_ref = f.get("cell", "")
            formula = f.get("formula", "")
            if cell_ref and formula:
                ws[cell_ref] = formula
                done.append(f"formula at {cell_ref}")

        # Named ranges
        for nr in named_ranges:
            name = nr.get("name", "")
            ref = nr.get("ref", "")
            if name and ref:
                from openpyxl.workbook.defined_name import DefinedName
                dn = DefinedName(name=name, attr_text=ref)
                wb.defined_names[name] = dn
                done.append(f"named range '{name}' = {ref}")

        # Conditional formatting
        for cf in cond_fmts:
            cf_range = cf.get("range", "")
            rule = cf.get("rule", "")
            color = cf.get("color", "FF0000")
            value = cf.get("value", 0)
            value2 = cf.get("value2", 0)

            if not cf_range:
                continue

            fill = PatternFill(start_color=color, end_color=color, fill_type="solid")

            if rule == "greater":
                ws.conditional_formatting.add(cf_range,
                    CellIsRule(operator="greaterThan", formula=[str(value)], fill=fill))
            elif rule == "less":
                ws.conditional_formatting.add(cf_range,
                    CellIsRule(operator="lessThan", formula=[str(value)], fill=fill))
            elif rule == "equal":
                ws.conditional_formatting.add(cf_range,
                    CellIsRule(operator="equal", formula=[str(value)], fill=fill))
            elif rule == "between":
                ws.conditional_formatting.add(cf_range,
                    CellIsRule(operator="between", formula=[str(value), str(value2)], fill=fill))
            elif rule == "color_scale":
                ws.conditional_formatting.add(cf_range, ColorScaleRule(
                    start_type="min", start_color="63BE7B",
                    mid_type="percentile", mid_value=50, mid_color="FFEB84",
                    end_type="max", end_color="F8696B",
                ))
            elif rule == "data_bar":
                ws.conditional_formatting.add(cf_range, DataBarRule(
                    start_type="min", start_value=0,
                    end_type="max", end_value=100,
                    color="638EC6",
                ))
            done.append(f"conditional format '{rule}' on {cf_range}")

        # Data validation
        from openpyxl.worksheet.datavalidation import DataValidation
        for dv_def in validations:
            dv_range = dv_def.get("range", "")
            dv_type = dv_def.get("type", "list")
            formula1 = dv_def.get("formula1", "")
            formula2 = dv_def.get("formula2", None)
            message = dv_def.get("message", "")

            if not dv_range or not formula1:
                continue

            # List validation: comma-separated → Excel list format
            if dv_type == "list" and "," in str(formula1) and not formula1.startswith("$"):
                formula1 = f'"{formula1}"'

            dv = DataValidation(type=dv_type, formula1=str(formula1),
                                formula2=str(formula2) if formula2 else None,
                                showErrorMessage=True)
            if message:
                dv.promptTitle = "Input"
                dv.prompt = message
            dv.add(dv_range)
            ws.add_data_validation(dv)
            done.append(f"validation '{dv_type}' on {dv_range}")

        wb.save(path)
        wb.close()
        return {"message": f"{len(done)} formula/format operations done", "details": done, "path": path}
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL — REPORT BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_report(args: Dict[str, Any]) -> Dict[str, Any]:
    """Build a complete professional Excel report in one shot.

    Combines data, summary stats, charts, and formatting into a polished file.

    Args:
        path: output .xlsx path
        title: report title
        data_sheet: main data definition:
            - name: sheet name (default "Data")
            - headers: column headers
            - data: rows of data
        summary_sheet: optional auto-generated summary:
            - name: sheet name (default "Summary")
            - metrics: list of {label, formula} e.g. {"label":"Total","formula":"=SUM(Data!B:B)"}
        chart: optional chart config (same as excel_chart args, applied to summary)
        style: "professional" | "colorful" | "minimal" (default professional)
        logo_text: optional text to put in top-left as branding
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    path = args.get("path", "")
    title = args.get("title", "Report")
    data_sheet_def = args.get("data_sheet", {})
    summary_def = args.get("summary_sheet", None)
    chart_def = args.get("chart", None)
    style_preset = args.get("style", "professional")
    logo_text = args.get("logo_text", "")

    if not path:
        return {"error": "path is required"}
    if not path.endswith(".xlsx"):
        path += ".xlsx"
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}

    THEME = {
        "professional": {
            "primary": "2F5496", "secondary": "D6E4F0",
            "header_font_color": "FFFFFF", "accent": "ED7D31",
        },
        "colorful": {
            "primary": "FF6B35", "secondary": "FFF0E8",
            "header_font_color": "FFFFFF", "accent": "4ECDC4",
        },
        "minimal": {
            "primary": "404040", "secondary": "F5F5F5",
            "header_font_color": "FFFFFF", "accent": "999999",
        },
    }
    theme = THEME.get(style_preset, THEME["professional"])

    try:
        _ensure_dir(path)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        data_name = data_sheet_def.get("name", "Data")
        headers = data_sheet_def.get("headers", [])
        data = data_sheet_def.get("data", [])

        # ── Summary sheet (first tab) ──
        if summary_def is not None or chart_def:
            sum_ws = wb.create_sheet(title=summary_def.get("name", "Summary") if summary_def else "Summary")

            # Title block
            sum_ws.merge_cells("A1:F1")
            tc = sum_ws["A1"]
            tc.value = title
            tc.font = Font(bold=True, size=18, color=theme["header_font_color"])
            tc.fill = PatternFill(start_color=theme["primary"], end_color=theme["primary"], fill_type="solid")
            tc.alignment = Alignment(horizontal="center", vertical="center")
            sum_ws.row_dimensions[1].height = 36

            if logo_text:
                sum_ws["A2"] = logo_text
                sum_ws["A2"].font = Font(size=9, color="999999", italic=True)

            # Date
            sum_ws["F2"] = f"Generated: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}"
            sum_ws["F2"].font = Font(size=9, color="999999")
            sum_ws["F2"].alignment = Alignment(horizontal="right")

            # Metrics
            if summary_def:
                metrics = summary_def.get("metrics", [])
                sum_ws["A4"] = "📊 Key Metrics"
                sum_ws["A4"].font = Font(bold=True, size=12, color=theme["primary"])
                for i, m in enumerate(metrics):
                    row = 5 + i
                    label_cell = sum_ws.cell(row=row, column=1, value=m.get("label", ""))
                    label_cell.font = Font(bold=True)
                    label_cell.fill = PatternFill(start_color=theme["secondary"], end_color=theme["secondary"], fill_type="solid")
                    val_cell = sum_ws.cell(row=row, column=2, value=m.get("formula", m.get("value", "")))
                    val_cell.font = Font(bold=True, color=theme["accent"])
                    val_cell.alignment = Alignment(horizontal="right")

        # ── Data sheet ──
        data_ws = wb.create_sheet(title=data_name)

        # Header row with title banner
        if title:
            data_ws.merge_cells(f"A1:{get_column_letter(max(len(headers), 1))}1")
            banner = data_ws["A1"]
            banner.value = title
            banner.font = Font(bold=True, size=14, color=theme["header_font_color"])
            banner.fill = PatternFill(start_color=theme["primary"], end_color=theme["primary"], fill_type="solid")
            banner.alignment = Alignment(horizontal="center")
            data_ws.row_dimensions[1].height = 28
            header_row = 2
        else:
            header_row = 1

        # Column headers
        side = Side(style="thin", color="B4C6E7")
        border = Border(left=side, right=side, top=side, bottom=side)

        for ci, h in enumerate(headers, 1):
            cell = data_ws.cell(row=header_row, column=ci, value=h)
            cell.font = Font(bold=True, color=theme["header_font_color"], size=11)
            cell.fill = PatternFill(start_color=theme["primary"], end_color=theme["primary"], fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border
        data_ws.row_dimensions[header_row].height = 22
        data_ws.freeze_panes = f"A{header_row + 1}"

        # Data rows with alternating colors
        alt_fill = PatternFill(start_color=theme["secondary"], end_color=theme["secondary"], fill_type="solid")
        for ri, row in enumerate(data):
            for ci, val in enumerate(row, 1):
                cell = data_ws.cell(row=header_row + 1 + ri, column=ci, value=_parse_value(val))
                cell.border = border
                if ri % 2 == 1:
                    cell.fill = alt_fill

        # Auto-width
        _auto_width(data_ws, headers, data)

        # Auto-filter
        if headers:
            last_col = get_column_letter(len(headers))
            last_row = header_row + len(data)
            data_ws.auto_filter.ref = f"A{header_row}:{last_col}{last_row}"

        # Totals row
        if headers and data:
            totals_row = header_row + len(data) + 1
            data_ws.cell(row=totals_row, column=1, value="TOTAL").font = Font(bold=True)
            for ci in range(2, len(headers) + 1):
                col_letter = get_column_letter(ci)
                # Only sum numeric-looking columns
                sample = [row[ci - 1] for row in data[:5] if ci <= len(row)]
                if any(isinstance(_parse_value(v), (int, float)) for v in sample):
                    cell = data_ws.cell(row=totals_row, column=ci,
                                        value=f"=SUM({col_letter}{header_row + 1}:{col_letter}{header_row + len(data)})")
                    cell.font = Font(bold=True, color=theme["accent"])
                    cell.number_format = "#,##0.##"

        wb.save(path)
        wb.close()

        return {
            "message": f"Report created: {path}",
            "path": path,
            "sheets": [s for s in ([summary_def.get("name","Summary") if summary_def else None, data_name]) if s],
            "data_rows": len(data),
            "style": style_preset,
        }
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL — PIVOT TABLE (via pandas)
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_pivot(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a pivot table from Excel data using pandas and write back to a new sheet.

    Args:
        path: source Excel or CSV file
        sheet: data sheet name (default: first)
        index: column name(s) for rows — string or list e.g. "Category" or ["Year","Category"]
        columns: column name for columns (optional)
        values: column name(s) to aggregate
        aggfunc: "sum" | "mean" | "count" | "max" | "min" (default: sum)
        output_sheet: sheet name for result (default: "Pivot")
        fill_na: value to fill NaN with (default 0)
    """
    import pandas as pd
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    path = args.get("path", "")
    sheet = args.get("sheet", 0)
    index = args.get("index", "")
    columns = args.get("columns", None)
    values = args.get("values", "")
    aggfunc = args.get("aggfunc", "sum")
    output_sheet = args.get("output_sheet", "Pivot")
    fill_na = args.get("fill_na", 0)

    if not path or not index or not values:
        return {"error": "path, index, and values are required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    AGGMAP = {"sum": "sum", "mean": "mean", "count": "count", "max": "max", "min": "min"}
    agg = AGGMAP.get(aggfunc, "sum")

    try:
        header_row = int(args.get("header_row", 0))  # 0-based row index for pandas

        ext = os.path.splitext(path)[1].lower()
        if ext == ".csv":
            df = pd.read_csv(path, header=header_row, skiprows=list(range(0, header_row)) if header_row > 0 else None)
        else:
            df = pd.read_excel(path, sheet_name=sheet if isinstance(sheet, int) else sheet, header=header_row)

        # Drop completely empty columns/rows
        df.dropna(how="all", axis=1, inplace=True)
        df.dropna(how="all", axis=0, inplace=True)
        df.reset_index(drop=True, inplace=True)

        # Normalize column names (strip whitespace)
        df.columns = [str(c).strip() for c in df.columns]
        index_cols = [i.strip() for i in (index if isinstance(index, list) else [index])]
        value_cols = [v.strip() for v in (values if isinstance(values, list) else [values])]
        cols_col = columns.strip() if columns else None

        pivot = pd.pivot_table(
            df,
            index=index_cols,
            columns=cols_col if cols_col else None,
            values=value_cols,
            aggfunc=agg,
            fill_value=fill_na,
        )
        pivot = pivot.reset_index()

        # Write pivot to new sheet in same workbook
        output_path = path if path.endswith(".xlsx") else os.path.splitext(path)[0] + "_pivot.xlsx"
        if os.path.isfile(output_path):
            wb = openpyxl.load_workbook(output_path)
            if output_sheet in wb.sheetnames:
                del wb[output_sheet]
        else:
            wb = openpyxl.Workbook()
            wb.remove(wb.active)

        ws = wb.create_sheet(title=output_sheet)

        # Write headers
        headers = list(pivot.columns.astype(str))
        for ci, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=ci, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
            cell.alignment = Alignment(horizontal="center")

        # Write data
        for ri, row in enumerate(pivot.itertuples(index=False), 2):
            for ci, val in enumerate(row, 1):
                ws.cell(row=ri, column=ci, value=val if not hasattr(val, 'item') else val.item())

        _auto_width(ws, headers, [list(r) for r in pivot.itertuples(index=False)])
        ws.freeze_panes = "A2"

        wb.save(output_path)
        wb.close()

        return {
            "message": f"Pivot table created in '{output_sheet}' sheet",
            "path": output_path,
            "rows": len(pivot),
            "columns": headers,
            "preview": pivot.head(10).to_string(),
        }
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL — VBA MACRO WRITER
# ═══════════════════════════════════════════════════════════════════════════════

def tool_excel_vba(args: Dict[str, Any]) -> Dict[str, Any]:
    """Write VBA macro code to an .xlsm file and optionally auto-open.

    Note: Creates/updates an .xlsm (macro-enabled) file. VBA runs when opened in Excel.

    Args:
        path: output .xlsm file path
        modules: list of VBA modules:
            - name: module name (e.g. "Module1", "AutoOpen")
            - code: VBA code string
        auto_run: procedure name to run automatically on open (e.g. "Main")
        description: human-readable description of what the macro does
    """
    # xlsm writing via openpyxl is limited — we use xlwings COM if available (macOS/Windows)
    # Fallback: write a .bas file + instructions

    path = args.get("path", "")
    modules = args.get("modules", [])
    auto_run = args.get("auto_run", "")
    description = args.get("description", "")

    if not path:
        return {"error": "path is required"}
    if not path.endswith(".xlsm"):
        path = os.path.splitext(path)[0] + ".xlsm"
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not modules:
        return {"error": "modules (list of {name, code}) required"}

    try:
        _ensure_dir(path)

        # Try xlwings approach (macOS with Excel)
        try:
            import xlwings as xw
            app = xw.App(visible=False)

            if os.path.isfile(path):
                wb = app.books.open(path)
            else:
                wb = app.books.add()

            for mod in modules:
                mod_name = mod.get("name", "Module1")
                code = mod.get("code", "")
                # Check if module exists
                try:
                    vba_mod = wb.api.VBProject.VBComponents(mod_name)
                    vba_mod.CodeModule.DeleteLines(1, vba_mod.CodeModule.CountOfLines)
                except Exception:
                    vba_mod = wb.api.VBProject.VBComponents.Add(1)  # 1 = vbext_ct_StdModule
                    vba_mod.Name = mod_name

                vba_mod.CodeModule.AddFromString(code)

            # Auto-run on open via ThisWorkbook
            if auto_run:
                try:
                    this_wb = wb.api.VBProject.VBComponents("ThisWorkbook")
                    open_code = f"Private Sub Workbook_Open()\n    {auto_run}\nEnd Sub"
                    this_wb.CodeModule.AddFromString(open_code)
                except Exception:
                    pass

            wb.save(path)
            wb.close()
            app.quit()

            return {
                "message": f"VBA macro written to {path}",
                "path": path,
                "modules": [m.get("name") for m in modules],
                "auto_run": auto_run,
            }

        except Exception as xlw_err:
            # Fallback: save .bas files + create companion .xlsx with instructions
            bas_files = []
            base_dir = os.path.dirname(path) or "/tmp"
            for mod in modules:
                mod_name = mod.get("name", "Module1")
                code = mod.get("code", "")
                bas_path = os.path.join(base_dir, f"{mod_name}.bas")
                with open(bas_path, "w", encoding="utf-8") as f:
                    f.write(f"Attribute VB_Name = \"{mod_name}\"\n")
                    f.write(code)
                bas_files.append(bas_path)

            # Create xlsx with instructions sheet
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "VBA Instructions"
            ws["A1"] = "VBA Macro Import Instructions"
            ws["A2"] = description or "Macro code saved as .bas files below"
            ws["A3"] = f"Auto-run procedure: {auto_run}" if auto_run else ""
            ws["A5"] = "Steps to import:"
            steps = [
                "1. Open Excel and press Alt+F11 to open VBA editor",
                "2. File → Import File → select each .bas file below",
                "3. Close VBA editor and save as .xlsm",
            ]
            for i, step in enumerate(steps):
                ws.cell(row=6 + i, column=1, value=step)
            ws["A10"] = "Generated .bas files:"
            for i, f in enumerate(bas_files):
                ws.cell(row=11 + i, column=1, value=f)

            xlsx_path = path.replace(".xlsm", "_vba_instructions.xlsx")
            wb.save(xlsx_path)
            wb.close()

            return {
                "message": "xlwings unavailable — VBA code saved as .bas files",
                "bas_files": bas_files,
                "instructions_xlsx": xlsx_path,
                "xlwings_error": str(xlw_err),
                "code_preview": {m.get("name"): m.get("code", "")[:200] for m in modules},
            }

    except Exception as e:
        return {"error": str(e)}


def tool_excel_vba_template(args: Dict[str, Any]) -> Dict[str, Any]:
    """Generate ready-to-use VBA code for common tasks.

    Args:
        template: one of:
            "format_all"      - Auto-format all sheets (headers, borders, autofit)
            "export_pdf"      - Export active sheet to PDF
            "send_email"      - Send workbook via Outlook
            "import_csv"      - Import CSV into a sheet
            "highlight_dups"  - Highlight duplicate values in selection
            "auto_sum"        - Add SUM row at bottom of all numeric columns
            "freeze_header"   - Freeze first row on all sheets
            "protect_sheets"  - Password protect all sheets
            "sort_data"       - Sort data by column A ascending
            "backup_save"     - Save backup copy with timestamp
        target_path: optional .xlsm path to write the generated macro to
    """
    TEMPLATES = {
        "format_all": {
            "name": "FormatAll",
            "code": """Sub FormatAll()
    Dim ws As Worksheet
    For Each ws In ThisWorkbook.Sheets
        ws.Cells.Font.Name = "Calibri"
        ws.Cells.Font.Size = 11
        With ws.Rows(1)
            .Font.Bold = True
            .Interior.Color = RGB(47, 84, 150)
            .Font.Color = RGB(255, 255, 255)
        End With
        ws.Cells.EntireColumn.AutoFit
        ws.Cells.Borders.LineStyle = xlContinuous
        ws.Cells.Borders.Weight = xlThin
    Next ws
    MsgBox "All sheets formatted!", vbInformation
End Sub""",
        },
        "export_pdf": {
            "name": "ExportPDF",
            "code": """Sub ExportPDF()
    Dim pdfPath As String
    pdfPath = ThisWorkbook.Path & "\\" & Left(ThisWorkbook.Name, Len(ThisWorkbook.Name) - 5) & "_" & Format(Now, "yyyymmdd_HHMM") & ".pdf"
    ActiveSheet.ExportAsFixedFormat Type:=xlTypePDF, Filename:=pdfPath, Quality:=xlQualityStandard
    MsgBox "PDF saved: " & pdfPath, vbInformation
End Sub""",
        },
        "send_email": {
            "name": "SendEmail",
            "code": """Sub SendEmail()
    Dim OutApp As Object
    Dim OutMail As Object
    Dim recipient As String
    recipient = InputBox("Enter recipient email:")
    If recipient = "" Then Exit Sub
    Set OutApp = CreateObject("Outlook.Application")
    Set OutMail = OutApp.CreateItem(0)
    With OutMail
        .To = recipient
        .Subject = "Report: " & ThisWorkbook.Name
        .Body = "Please find attached the report."
        .Attachments.Add ThisWorkbook.FullName
        .Send
    End With
    MsgBox "Email sent to " & recipient, vbInformation
End Sub""",
        },
        "import_csv": {
            "name": "ImportCSV",
            "code": """Sub ImportCSV()
    Dim csvPath As String
    Dim ws As Worksheet
    csvPath = Application.GetOpenFilename("CSV Files (*.csv), *.csv")
    If csvPath = "False" Then Exit Sub
    Set ws = ThisWorkbook.Sheets.Add(After:=ThisWorkbook.Sheets(ThisWorkbook.Sheets.Count))
    ws.Name = "Imported_" & Format(Now, "HHMMSS")
    With ws.QueryTables.Add(Connection:="TEXT;" & csvPath, Destination:=ws.Range("A1"))
        .TextFileCommaDelimiter = True
        .Refresh BackgroundQuery:=False
    End With
    MsgBox "CSV imported to sheet '" & ws.Name & "'", vbInformation
End Sub""",
        },
        "highlight_dups": {
            "name": "HighlightDuplicates",
            "code": """Sub HighlightDuplicates()
    Dim rng As Range
    Dim cell As Range
    Dim dict As Object
    Set rng = Selection
    Set dict = CreateObject("Scripting.Dictionary")
    ' First pass: count
    For Each cell In rng
        If cell.Value <> "" Then
            If dict.Exists(CStr(cell.Value)) Then
                dict(CStr(cell.Value)) = dict(CStr(cell.Value)) + 1
            Else
                dict.Add CStr(cell.Value), 1
            End If
        End If
    Next cell
    ' Second pass: highlight
    For Each cell In rng
        If dict.Exists(CStr(cell.Value)) And dict(CStr(cell.Value)) > 1 Then
            cell.Interior.Color = RGB(255, 199, 206)
        End If
    Next cell
    MsgBox "Duplicates highlighted in red.", vbInformation
End Sub""",
        },
        "auto_sum": {
            "name": "AutoSum",
            "code": """Sub AutoSum()
    Dim ws As Worksheet
    Dim lastRow As Long
    Dim lastCol As Long
    Dim c As Integer
    Set ws = ActiveSheet
    lastRow = ws.Cells(ws.Rows.Count, 1).End(xlUp).Row
    lastCol = ws.Cells(1, ws.Columns.Count).End(xlToLeft).Column
    ws.Cells(lastRow + 1, 1).Value = "TOTAL"
    ws.Cells(lastRow + 1, 1).Font.Bold = True
    For c = 2 To lastCol
        If IsNumeric(ws.Cells(2, c).Value) Then
            ws.Cells(lastRow + 1, c).Formula = "=SUM(" & ws.Cells(2, c).Address(False, False) & ":" & ws.Cells(lastRow, c).Address(False, False) & ")"
            ws.Cells(lastRow + 1, c).Font.Bold = True
        End If
    Next c
End Sub""",
        },
        "freeze_header": {
            "name": "FreezeHeaders",
            "code": """Sub FreezeHeaders()
    Dim ws As Worksheet
    For Each ws In ThisWorkbook.Sheets
        ws.Activate
        ws.Rows(2).Select
        ActiveWindow.FreezePanes = True
    Next ws
    ThisWorkbook.Sheets(1).Activate
    MsgBox "Headers frozen on all sheets.", vbInformation
End Sub""",
        },
        "protect_sheets": {
            "name": "ProtectSheets",
            "code": """Sub ProtectSheets()
    Dim ws As Worksheet
    Dim pwd As String
    pwd = InputBox("Enter protection password (leave blank for no password):")
    For Each ws In ThisWorkbook.Sheets
        ws.Protect Password:=pwd, DrawingObjects:=True, Contents:=True, Scenarios:=True
    Next ws
    MsgBox "All sheets protected.", vbInformation
End Sub""",
        },
        "sort_data": {
            "name": "SortData",
            "code": """Sub SortData()
    Dim ws As Worksheet
    Dim lastRow As Long
    Dim lastCol As Long
    Set ws = ActiveSheet
    lastRow = ws.Cells(ws.Rows.Count, 1).End(xlUp).Row
    lastCol = ws.Cells(1, ws.Columns.Count).End(xlToLeft).Column
    ws.Range(ws.Cells(1, 1), ws.Cells(lastRow, lastCol)).Sort _
        Key1:=ws.Columns(1), Order1:=xlAscending, Header:=xlYes
    MsgBox "Data sorted by column A (ascending).", vbInformation
End Sub""",
        },
        "backup_save": {
            "name": "BackupSave",
            "code": """Sub BackupSave()
    Dim backupPath As String
    Dim baseName As String
    baseName = Left(ThisWorkbook.Name, InStrRev(ThisWorkbook.Name, ".") - 1)
    backupPath = ThisWorkbook.Path & "\\" & baseName & "_backup_" & Format(Now, "yyyymmdd_HHMM") & ".xlsx"
    ThisWorkbook.SaveCopyAs Filename:=backupPath
    MsgBox "Backup saved: " & backupPath, vbInformation
End Sub""",
        },
    }

    template = args.get("template", "")
    target_path = args.get("target_path", "")

    if not template or template not in TEMPLATES:
        return {
            "available_templates": list(TEMPLATES.keys()),
            "error": f"Unknown template '{template}'. Choose from above.",
        }

    tpl = TEMPLATES[template]
    result = {
        "template": template,
        "module_name": tpl["name"],
        "code": tpl["code"],
        "usage": f"Call this macro via: Alt+F8 → select '{tpl['name']}' → Run",
    }

    if target_path:
        write_result = tool_excel_vba({
            "path": target_path,
            "modules": [{"name": tpl["name"], "code": tpl["code"]}],
        })
        result["write_result"] = write_result

    return result


# ═══════════════════════════════════════════════════════════════════════════════
#  TOOL REGISTRY
# ═══════════════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════════════
#  PDF TOOLS — Read, extract tables, analyze financial reports
# ═══════════════════════════════════════════════════════════════════════════════

def tool_pdf_read(args: Dict[str, Any]) -> Dict[str, Any]:
    """Read a PDF file. Extract text and metadata.

    Args:
        path: PDF file path
        pages: page range to read (default: all). e.g. "1-5", "1,3,5", "all"
        max_chars: max characters to return (default 8000)
    """
    try:
        import pdfplumber
    except ImportError:
        return {"error": "pdfplumber not installed. Run: pip install pdfplumber"}

    path = args.get("path", "")
    pages_arg = str(args.get("pages", "all"))
    max_chars = int(args.get("max_chars", 8000))

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)

            # Parse page range
            if pages_arg == "all":
                page_indices = list(range(total_pages))
            else:
                page_indices = []
                for part in pages_arg.replace(" ", "").split(","):
                    if "-" in part:
                        start, end = part.split("-", 1)
                        page_indices.extend(range(int(start) - 1, min(int(end), total_pages)))
                    else:
                        idx = int(part) - 1
                        if 0 <= idx < total_pages:
                            page_indices.append(idx)

            # Extract text
            text_parts = []
            for i in page_indices:
                page = pdf.pages[i]
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

            full_text = "\n\n".join(text_parts)
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + "\n...[truncated]"

            return {
                "text": full_text,
                "total_pages": total_pages,
                "pages_read": len(page_indices),
                "file": os.path.basename(path),
                "file_size_kb": round(os.path.getsize(path) / 1024, 1),
            }
    except Exception as e:
        return {"error": f"Failed to read PDF: {str(e)}"}


def tool_pdf_tables(args: Dict[str, Any]) -> Dict[str, Any]:
    """Extract tables from PDF. Essential for financial reports with structured data.

    Args:
        path: PDF file path
        pages: page range (default: "all"). e.g. "1-5", "1,3"
        table_index: specific table index to return (default: all tables)
        format: "rows" (list of lists) or "dict" (list of dicts with headers) — default "dict"
    """
    try:
        import pdfplumber
    except ImportError:
        return {"error": "pdfplumber not installed. Run: pip install pdfplumber"}

    path = args.get("path", "")
    pages_arg = str(args.get("pages", "all"))
    table_index = args.get("table_index", None)
    fmt = args.get("format", "dict")

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)

            # Parse page range
            if pages_arg == "all":
                page_indices = list(range(total_pages))
            else:
                page_indices = []
                for part in pages_arg.replace(" ", "").split(","):
                    if "-" in part:
                        start, end = part.split("-", 1)
                        page_indices.extend(range(int(start) - 1, min(int(end), total_pages)))
                    else:
                        idx = int(part) - 1
                        if 0 <= idx < total_pages:
                            page_indices.append(idx)

            all_tables = []
            for i in page_indices:
                page = pdf.pages[i]
                tables = page.extract_tables()
                for t_idx, table in enumerate(tables):
                    if not table or len(table) < 2:
                        continue

                    # Clean cells
                    cleaned = []
                    for row in table:
                        cleaned.append([
                            (cell.strip() if isinstance(cell, str) else str(cell) if cell is not None else "")
                            for cell in row
                        ])

                    if fmt == "dict" and cleaned:
                        headers = cleaned[0]
                        # Ensure unique headers
                        seen = {}
                        unique_headers = []
                        for h in headers:
                            h = h or f"col_{len(unique_headers)}"
                            if h in seen:
                                seen[h] += 1
                                h = f"{h}_{seen[h]}"
                            else:
                                seen[h] = 0
                            unique_headers.append(h)

                        dict_rows = []
                        for row in cleaned[1:]:
                            row_dict = {}
                            for j, val in enumerate(row):
                                if j < len(unique_headers):
                                    row_dict[unique_headers[j]] = val
                            dict_rows.append(row_dict)

                        all_tables.append({
                            "page": i + 1,
                            "table_index": t_idx,
                            "headers": unique_headers,
                            "rows": dict_rows[:100],  # cap at 100 rows
                            "total_rows": len(dict_rows),
                        })
                    else:
                        all_tables.append({
                            "page": i + 1,
                            "table_index": t_idx,
                            "rows": cleaned[:100],
                            "total_rows": len(cleaned),
                        })

            # Filter by table_index if specified
            if table_index is not None:
                table_index = int(table_index)
                if 0 <= table_index < len(all_tables):
                    all_tables = [all_tables[table_index]]

            return {
                "tables": all_tables,
                "total_tables": len(all_tables),
                "file": os.path.basename(path),
            }
    except Exception as e:
        return {"error": f"Failed to extract tables: {str(e)}"}


def tool_pdf_financial(args: Dict[str, Any]) -> Dict[str, Any]:
    """Analyze a financial report PDF. Extracts key metrics, income statement, balance sheet data.

    This tool combines text extraction + table extraction + financial number parsing.
    Best for: annual reports, quarterly earnings, balance sheets, income statements.

    Args:
        path: PDF file path
        pages: page range (default: "all")
        focus: "overview" | "income" | "balance" | "cashflow" | "ratios" | "all" (default: "all")
    """
    try:
        import pdfplumber
    except ImportError:
        return {"error": "pdfplumber not installed. Run: pip install pdfplumber"}
    import re

    path = args.get("path", "")
    pages_arg = str(args.get("pages", "all"))
    focus = args.get("focus", "all")

    if not path:
        return {"error": "path is required"}
    if not _check_path(path):
        return {"error": f"Path not allowed: {path}"}
    if not os.path.isfile(path):
        return {"error": f"File not found: {path}"}

    def _parse_number(s: str) -> Optional[float]:
        """Parse a financial number string to float. Handles (parens), commas, %, B/M/K."""
        if not s or not isinstance(s, str):
            return None
        s = s.strip()
        if not s:
            return None

        negative = False
        if s.startswith("(") and s.endswith(")"):
            negative = True
            s = s[1:-1]
        if s.startswith("-"):
            negative = True
            s = s[1:]

        # Remove currency symbols and spaces
        s = re.sub(r'[$ €£¥₫VNĐđ\s]', '', s)

        # Handle suffixes
        multiplier = 1
        s_upper = s.upper()
        if s_upper.endswith('B'):
            multiplier = 1_000_000_000
            s = s[:-1]
        elif s_upper.endswith('M') or s_upper.endswith('TR'):
            multiplier = 1_000_000
            s = s[:-1] if s_upper.endswith('M') else s[:-2]
        elif s_upper.endswith('K') or s_upper.endswith('N'):
            multiplier = 1_000
            s = s[:-1]
        elif s_upper.endswith('%'):
            s = s[:-1]
            multiplier = 0.01

        # Handle Vietnamese number format (dot as thousands separator, comma as decimal)
        if '.' in s and ',' in s:
            # Both present: dots are thousands, comma is decimal
            s = s.replace('.', '').replace(',', '.')
        elif ',' in s:
            # Check if comma is decimal or thousands
            parts = s.split(',')
            if len(parts) == 2 and len(parts[1]) <= 2:
                s = s.replace(',', '.')  # decimal
            else:
                s = s.replace(',', '')  # thousands
        elif '.' in s:
            parts = s.split('.')
            if len(parts) == 2 and len(parts[1]) == 3:
                s = s.replace('.', '')  # thousands separator

        try:
            val = float(s) * multiplier
            return -val if negative else val
        except ValueError:
            return None

    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)

            if pages_arg == "all":
                page_indices = list(range(total_pages))
            else:
                page_indices = []
                for part in pages_arg.replace(" ", "").split(","):
                    if "-" in part:
                        start, end = part.split("-", 1)
                        page_indices.extend(range(int(start) - 1, min(int(end), total_pages)))
                    else:
                        idx = int(part) - 1
                        if 0 <= idx < total_pages:
                            page_indices.append(idx)

            # Extract all text and tables
            all_text = ""
            all_tables = []
            for i in page_indices:
                page = pdf.pages[i]
                text = page.extract_text() or ""
                all_text += f"\n{text}"

                for table in page.extract_tables():
                    if table and len(table) >= 2:
                        cleaned = []
                        for row in table:
                            cleaned.append([
                                (cell.strip() if isinstance(cell, str) else str(cell) if cell is not None else "")
                                for cell in row
                            ])
                        all_tables.append({"page": i + 1, "data": cleaned})

            # ── Financial keyword detection ──
            text_lower = all_text.lower()

            # Detect report language
            is_vietnamese = any(kw in text_lower for kw in ["doanh thu", "lợi nhuận", "tài sản", "nguồn vốn", "bảng cân đối"])

            # Keywords for section detection
            income_keywords = (
                ["doanh thu", "lợi nhuận", "chi phí", "thu nhập", "kết quả kinh doanh"]
                if is_vietnamese else
                ["revenue", "income", "profit", "expense", "earnings", "net income", "gross profit", "operating"]
            )
            balance_keywords = (
                ["tài sản", "nợ phải trả", "vốn chủ", "nguồn vốn", "bảng cân đối"]
                if is_vietnamese else
                ["assets", "liabilities", "equity", "balance sheet", "total assets", "current assets"]
            )
            cashflow_keywords = (
                ["lưu chuyển tiền tệ", "dòng tiền", "tiền thuần"]
                if is_vietnamese else
                ["cash flow", "operating activities", "investing", "financing", "free cash flow"]
            )

            # ── Parse financial tables ──
            financial_data = {
                "income_statement": [],
                "balance_sheet": [],
                "cash_flow": [],
                "other_tables": [],
            }

            for tbl in all_tables:
                data = tbl["data"]
                # Flatten table text to detect category
                flat = " ".join(" ".join(row) for row in data).lower()

                categorized = False
                if any(kw in flat for kw in income_keywords):
                    financial_data["income_statement"].append(tbl)
                    categorized = True
                if any(kw in flat for kw in balance_keywords):
                    financial_data["balance_sheet"].append(tbl)
                    categorized = True
                if any(kw in flat for kw in cashflow_keywords):
                    financial_data["cash_flow"].append(tbl)
                    categorized = True
                if not categorized:
                    financial_data["other_tables"].append(tbl)

            # ── Extract key metrics from tables ──
            key_metrics = {}
            metric_patterns = {
                # Vietnamese
                "doanh_thu": r"(?:doanh thu|tổng doanh thu|net revenue|revenue)",
                "loi_nhuan_gop": r"(?:lợi nhuận gộp|gross profit)",
                "loi_nhuan_rong": r"(?:lợi nhuận.*sau thuế|lợi nhuận ròng|net (?:income|profit))",
                "tong_tai_san": r"(?:tổng (?:cộng )?tài sản|total assets)",
                "von_chu_so_huu": r"(?:vốn chủ sở hữu|(?:total )?(?:shareholders['\s]?)?equity)",
                "no_phai_tra": r"(?:tổng nợ|nợ phải trả|total (?:liabilities|debt))",
                "eps": r"(?:eps|earning.{0,5}per.{0,5}share|lãi.*cổ phiếu)",
            }

            for tbl in all_tables:
                for row in tbl["data"]:
                    if not row:
                        continue
                    label = (row[0] or "").strip()
                    label_lower = label.lower()

                    for metric_key, pattern in metric_patterns.items():
                        if metric_key not in key_metrics and re.search(pattern, label_lower):
                            # Find the first numeric value in the row
                            for cell in row[1:]:
                                val = _parse_number(cell)
                                if val is not None:
                                    key_metrics[metric_key] = {
                                        "label": label,
                                        "value": val,
                                        "raw": cell,
                                    }
                                    break

            # ── Calculate ratios ──
            ratios = {}
            if "doanh_thu" in key_metrics and "loi_nhuan_gop" in key_metrics:
                rev = key_metrics["doanh_thu"]["value"]
                gp = key_metrics["loi_nhuan_gop"]["value"]
                if rev and rev != 0:
                    ratios["gross_margin"] = f"{(gp / rev * 100):.1f}%"

            if "doanh_thu" in key_metrics and "loi_nhuan_rong" in key_metrics:
                rev = key_metrics["doanh_thu"]["value"]
                np_ = key_metrics["loi_nhuan_rong"]["value"]
                if rev and rev != 0:
                    ratios["net_margin"] = f"{(np_ / rev * 100):.1f}%"

            if "loi_nhuan_rong" in key_metrics and "von_chu_so_huu" in key_metrics:
                np_ = key_metrics["loi_nhuan_rong"]["value"]
                eq = key_metrics["von_chu_so_huu"]["value"]
                if eq and eq != 0:
                    ratios["roe"] = f"{(np_ / eq * 100):.1f}%"

            if "loi_nhuan_rong" in key_metrics and "tong_tai_san" in key_metrics:
                np_ = key_metrics["loi_nhuan_rong"]["value"]
                ta = key_metrics["tong_tai_san"]["value"]
                if ta and ta != 0:
                    ratios["roa"] = f"{(np_ / ta * 100):.1f}%"

            if "no_phai_tra" in key_metrics and "von_chu_so_huu" in key_metrics:
                debt = key_metrics["no_phai_tra"]["value"]
                eq = key_metrics["von_chu_so_huu"]["value"]
                if eq and eq != 0:
                    ratios["debt_to_equity"] = f"{(debt / eq):.2f}"

            # ── Build result ──
            result = {
                "file": os.path.basename(path),
                "total_pages": total_pages,
                "pages_analyzed": len(page_indices),
                "language": "Vietnamese" if is_vietnamese else "English",
                "tables_found": len(all_tables),
            }

            if focus in ("all", "overview"):
                result["key_metrics"] = key_metrics
                result["ratios"] = ratios

            if focus in ("all", "income"):
                result["income_tables"] = len(financial_data["income_statement"])
                for i, tbl in enumerate(financial_data["income_statement"][:2]):
                    result[f"income_table_{i}"] = {
                        "page": tbl["page"],
                        "data": tbl["data"][:30],
                    }

            if focus in ("all", "balance"):
                result["balance_tables"] = len(financial_data["balance_sheet"])
                for i, tbl in enumerate(financial_data["balance_sheet"][:2]):
                    result[f"balance_table_{i}"] = {
                        "page": tbl["page"],
                        "data": tbl["data"][:30],
                    }

            if focus in ("all", "cashflow"):
                result["cashflow_tables"] = len(financial_data["cash_flow"])
                for i, tbl in enumerate(financial_data["cash_flow"][:2]):
                    result[f"cashflow_table_{i}"] = {
                        "page": tbl["page"],
                        "data": tbl["data"][:30],
                    }

            if focus == "ratios":
                result["key_metrics"] = key_metrics
                result["ratios"] = ratios

            # Add text snippet for context
            if len(all_text) > 200:
                result["text_preview"] = all_text[:1500] + "..."
            else:
                result["text_preview"] = all_text

            return result

    except Exception as e:
        return {"error": f"Financial analysis failed: {str(e)}"}


OFFICE_TOOLS = {
    # Excel core
    "excel_read": tool_excel_read,
    "excel_write": tool_excel_write,
    "excel_edit": tool_excel_edit,
    "excel_chart": tool_excel_chart,
    "excel_format": tool_excel_format,
    "excel_analyze": tool_excel_analyze,
    # Excel advanced
    "excel_rows_cols": tool_excel_rows_cols,
    "excel_sheets": tool_excel_sheets,
    "excel_formula": tool_excel_formula,
    "excel_report": tool_excel_report,
    "excel_pivot": tool_excel_pivot,
    "excel_vba": tool_excel_vba,
    "excel_vba_template": tool_excel_vba_template,
    # Word
    "word_read": tool_word_read,
    "word_write": tool_word_write,
    # PowerPoint
    "pptx_read": tool_pptx_read,
    "pptx_write": tool_pptx_write,
    # Conversion
    "csv_to_excel": tool_csv_to_excel,
    # PDF
    "pdf_read": tool_pdf_read,
    "pdf_tables": tool_pdf_tables,
    "pdf_financial": tool_pdf_financial,
}

OFFICE_TOOL_DEFINITIONS = """
OFFICE: excel_read{path,sheet?} | excel_write{path,headers,data,style?} | excel_edit{path,edits:[{cell,value}]} | excel_rows_cols{path,operations} | excel_sheets{path,operations} | excel_formula{path,formulas} | excel_chart{path,chart_type,title,data_range} | excel_format{path,formats} | excel_analyze{path,column,operation} | excel_report{path,title,data_sheet} | excel_pivot{path,index,values} | excel_vba{path,modules} | word_read{path} | word_write{path,content} | pptx_read{path} | pptx_write{path,title,slides} | csv_to_excel{input_path}
PDF: pdf_read{path,pages?} → extract text | pdf_tables{path,pages?,format?} → extract tables as rows/dicts | pdf_financial{path,pages?,focus?} → analyze financial report (income/balance/cashflow/ratios). focus: overview|income|balance|cashflow|ratios|all
"""
